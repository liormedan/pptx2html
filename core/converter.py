from pathlib import Path
from pptx import Presentation
from .styles import StyleGenerator
from .animations import AnimationHandler
from .tables import TableHandler
import base64
from io import BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
import math
import logging
from jinja2 import Template
import os

logger = logging.getLogger(__name__)

class PowerPointConverter:
    def __init__(self):
        """אתחול הממיר"""
        self.style_generator = StyleGenerator()
        self.animation_handler = AnimationHandler()
        self.table_handler = TableHandler()
        self.animation_step = 0

    def _add_animation_to_element(self, element_id, shape_index, slide_index):
        """הוספת אנימציה לאלמנט"""
        animation = self.animation_handler.create_animation(
            effect_type='FADE',  # ברירת מחדל
            step=shape_index,    # כל צורה תופיע בשלב נפרד
            delay=0.2 * shape_index,  # דיליי הדרגתי
            duration=1.0
        )
        return animation

    def convert(self, pptx_path, output_path=None):
        """
        ממיר מצגת PowerPoint לקובץ HTML
        
        Args:
            pptx_path (str): נתיב לקובץ PowerPoint
            output_path (str, optional): נתיב לשמירת קובץ ה-HTML. אם לא צוין, ישמר באותה תיקייה
        
        Returns:
            str: נתיב לקובץ ה-HTML שנוצר
        """
        try:
            presentation = Presentation(pptx_path)
            html_content = self._convert_presentation(presentation)
            
            if output_path is None:
                # שמירה באותה תיקייה כמו קובץ המקור
                output_path = str(Path(pptx_path).with_suffix('.html'))
                
            # שמירת הקובץ
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
                
            return output_path
            
        except Exception as e:
            raise Exception(f"שגיאה בהמרת המצגת: {str(e)}")
        
    def _convert_presentation(self, presentation):
        """ממיר את המצגת כולה ל-HTML"""
        slides_html = []
        for i, slide in enumerate(presentation.slides):
            slide_html = self._convert_slide(slide, i)
            slides_html.append(slide_html)

        # יצירת ה-HTML המלא
        html = self._generate_html_content(slides_html)
        return html
    
    def _generate_html_content(self, slides_data):
        """יוצר את תוכן ה-HTML עם כל התכונות המתקדמות"""
        html_template = """
<!DOCTYPE html>
<html dir="rtl" lang="he">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>מצגת מומרת</title>
    <style>
        :root {
            --primary-color: #2196F3;
            --background-color: #ffffff;
            --text-color: #333333;
            --thumbnail-size: 120px;
        }

        [data-theme="dark"] {
            --background-color: #1a1a1a;
            --text-color: #ffffff;
        }

        body {
            margin: 0;
            padding: 0;
            background-color: var(--background-color);
            color: var(--text-color);
            font-family: system-ui, -apple-system, sans-serif;
            transition: background-color 0.3s, color 0.3s;
            overflow-x: hidden;
        }

        #presentation-container {
            display: flex;
            height: 100vh;
            max-width: 100vw;
        }

        #slides-container {
            flex: 1;
            overflow: hidden;
            position: relative;
            display: flex;
            justify-content: center;
            align-items: center;
        }

        .slide {
            position: absolute;
            width: 100%;
            height: 100%;
            opacity: 0;
            transition: opacity 0.5s;
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 20px;
            box-sizing: border-box;
        }

        .slide.active {
            opacity: 1;
            z-index: 1;
        }

        /* תיקון גודל תמונות */
        .slide img {
            max-width: 100%;
            max-height: 80vh;
            object-fit: contain;
            margin: auto;
        }

        .slide-content {
            max-width: 90%;
            max-height: 90vh;
            margin: auto;
            overflow: auto;
            display: flex;
            flex-direction: column;
            align-items: center;
        }

        /* שאר הסגנונות נשארים אותו דבר... */
        
        #export-panel {
            position: fixed;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            background: var(--background-color);
            padding: 20px;
            border-radius: 10px;
            box-shadow: 0 0 20px rgba(0,0,0,0.3);
            z-index: 1000;
            display: none;
        }

        #export-panel textarea {
            width: 100%;
            min-height: 150px;
            margin: 10px 0;
            padding: 10px;
            direction: ltr;
            font-family: monospace;
        }

        .close-button {
            position: absolute;
            top: 10px;
            right: 10px;
            background: none;
            border: none;
            font-size: 20px;
            cursor: pointer;
        }

        /* סגנונות נוספים נשארים כמו שהם... */
        
        #thumbnails-panel {
            width: var(--thumbnail-size);
            background: rgba(0,0,0,0.1);
            overflow-y: auto;
            padding: 10px;
            display: none;
        }

        .thumbnail {
            width: 100px;
            height: 75px;
            margin-bottom: 10px;
            cursor: pointer;
            border: 2px solid transparent;
            background-size: cover;
            background-position: center;
        }

        .thumbnail.active {
            border-color: var(--primary-color);
        }

        #controls {
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            z-index: 100;
            display: flex;
            gap: 10px;
            background: rgba(0,0,0,0.5);
            padding: 10px;
            border-radius: 20px;
        }

        .control-button {
            background: transparent;
            border: none;
            color: white;
            cursor: pointer;
            padding: 5px 10px;
            font-size: 16px;
            border-radius: 5px;
        }

        .control-button:hover {
            background: rgba(255,255,255,0.1);
        }

        #settings-panel {
            position: fixed;
            top: 20px;
            right: 20px;
            background: var(--background-color);
            padding: 15px;
            border-radius: 10px;
            box-shadow: 0 0 10px rgba(0,0,0,0.2);
            z-index: 100;
            display: none;
        }
    </style>
</head>
<body>
    <div id="presentation-container">
        <div id="thumbnails-panel"></div>
        <div id="slides-container">
            {% for slide in slides %}
            <div class="slide" id="slide-{{ loop.index }}">
                {{ slide.content | safe }}
            </div>
            {% endfor %}
        </div>
    </div>

    <div id="controls">
        <button class="control-button" onclick="previousSlide()">&#10094;</button>
        <button class="control-button" onclick="toggleThumbnails()">תמונות ממוזערות</button>
        <button class="control-button" onclick="toggleFullscreen()">מסך מלא</button>
        <button class="control-button" onclick="toggleSettings()">הגדרות</button>
        <button class="control-button" onclick="exportToWix()">ייצוא ל-Wix</button>
        <button class="control-button" onclick="nextSlide()">&#10095;</button>
    </div>

    <div id="settings-panel">
        <div class="settings-row">
            <label>מצב כהה</label>
            <input type="checkbox" id="dark-mode" onchange="toggleTheme()">
        </div>
        <div class="settings-row">
            <label>גודל טקסט</label>
            <select id="font-size" onchange="changeFontSize()">
                <option value="small">קטן</option>
                <option value="medium" selected>בינוני</option>
                <option value="large">גדול</option>
            </select>
        </div>
    </div>

    <div id="export-panel">
        <button class="close-button" onclick="closeExportPanel()">&times;</button>
        <h3>קוד להטמעה ב-Wix</h3>
        <p>העתק את הקוד הבא והדבק אותו ברכיב HTML מותאם אישית ב-Wix:</p>
        <textarea id="embed-code" readonly></textarea>
        <button class="control-button" onclick="copyEmbedCode()">העתק קוד</button>
    </div>

    <script>
        let currentSlide = 1;
        const totalSlides = {{ slides|length }};
        
        // טעינת העדפות משתמש
        document.addEventListener('DOMContentLoaded', () => {
            loadUserPreferences();
            showSlide(currentSlide);
            createThumbnails();
        });

        function showSlide(n) {
            const slides = document.querySelectorAll('.slide');
            const thumbnails = document.querySelectorAll('.thumbnail');
            
            currentSlide = n;
            if (currentSlide > totalSlides) currentSlide = 1;
            if (currentSlide < 1) currentSlide = totalSlides;
            
            slides.forEach(slide => slide.classList.remove('active'));
            thumbnails.forEach(thumb => thumb.classList.remove('active'));
            
            slides[currentSlide - 1].classList.add('active');
            if (thumbnails[currentSlide - 1]) {
                thumbnails[currentSlide - 1].classList.add('active');
                thumbnails[currentSlide - 1].scrollIntoView({ behavior: 'smooth', block: 'nearest' });
            }
        }

        function nextSlide() {
            showSlide(currentSlide + 1);
        }

        function previousSlide() {
            showSlide(currentSlide - 1);
        }

        function createThumbnails() {
            const thumbnailsPanel = document.getElementById('thumbnails-panel');
            const slides = document.querySelectorAll('.slide');
            
            slides.forEach((slide, index) => {
                const thumbnail = document.createElement('div');
                thumbnail.className = 'thumbnail';
                thumbnail.onclick = () => showSlide(index + 1);
                thumbnail.style.backgroundImage = `url('data:image/svg+xml,<svg xmlns="http://www.w3.org/2000/svg" width="100" height="75"><rect width="100%" height="100%" fill="%23f0f0f0"/><text x="50%" y="50%" dominant-baseline="middle" text-anchor="middle" font-family="system-ui" font-size="20" fill="%23666">${index + 1}</text></svg>')`;
                thumbnailsPanel.appendChild(thumbnail);
            });
        }

        function toggleThumbnails() {
            const panel = document.getElementById('thumbnails-panel');
            panel.style.display = panel.style.display === 'none' ? 'block' : 'none';
        }

        function toggleFullscreen() {
            if (!document.fullscreenElement) {
                document.documentElement.requestFullscreen();
            } else {
                document.exitFullscreen();
            }
        }

        function toggleSettings() {
            const panel = document.getElementById('settings-panel');
            panel.style.display = panel.style.display === 'none' ? 'block' : 'none';
        }

        function toggleTheme() {
            const isDark = document.getElementById('dark-mode').checked;
            document.documentElement.setAttribute('data-theme', isDark ? 'dark' : 'light');
            localStorage.setItem('theme', isDark ? 'dark' : 'light');
        }

        function changeFontSize() {
            const size = document.getElementById('font-size').value;
            const sizes = {
                small: '14px',
                medium: '16px',
                large: '18px'
            };
            document.body.style.fontSize = sizes[size];
            localStorage.setItem('fontSize', size);
        }

        function loadUserPreferences() {
            // טעינת ערכת נושא
            const theme = localStorage.getItem('theme') || 'light';
            document.getElementById('dark-mode').checked = theme === 'dark';
            document.documentElement.setAttribute('data-theme', theme);

            // טעינת גודל טקסט
            const fontSize = localStorage.getItem('fontSize') || 'medium';
            document.getElementById('font-size').value = fontSize;
            changeFontSize();
        }

        // תמיכה בניווט מקלדת
        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight') previousSlide();
            else if (e.key === 'ArrowLeft') nextSlide();
            else if (e.key === 'f') toggleFullscreen();
        });

        // תמיכה במחוות מגע
        let touchStartX = 0;
        document.addEventListener('touchstart', e => {
            touchStartX = e.touches[0].clientX;
        });

        document.addEventListener('touchend', e => {
            const touchEndX = e.changedTouches[0].clientX;
            const diff = touchStartX - touchEndX;
            
            if (Math.abs(diff) > 50) {
                if (diff > 0) nextSlide();
                else previousSlide();
            }
        });

        function exportToWix() {
            const panel = document.getElementById('export-panel');
            const textarea = document.getElementById('embed-code');
            const currentUrl = window.location.href;
            
            const embedCode = `
<!-- קוד להטמעת המצגת -->
<div class="presentation-container" style="position: relative; width: 100%;">
    <iframe 
        src="${currentUrl}"
        style="width: 100%; height: 600px; border: none; overflow: hidden;"
        allowfullscreen="true"
        loading="lazy"
    ></iframe>
</div>

<!-- סקריפט להתאמה אוטומטית של גובה -->
<script>
window.addEventListener('message', function(e) {
    if (e.data && e.data.type === 'resize') {
        const iframe = document.querySelector('.presentation-container iframe');
        if (iframe) {
            iframe.style.height = e.data.height + 'px';
        }
    }
});
<\/script>`;
            
            textarea.value = embedCode;
            panel.style.display = 'block';
        }

        function closeExportPanel() {
            document.getElementById('export-panel').style.display = 'none';
        }

        function copyEmbedCode() {
            const textarea = document.getElementById('embed-code');
            textarea.select();
            document.execCommand('copy');
            alert('הקוד הועתק ללוח!');
        }
    </script>
</body>
</html>
"""
        template = Template(html_template)
        return template.render(slides=slides_data)

    def _convert_slide(self, slide, index):
        """ממיר שקופית בודדת ל-HTML"""
        try:
            slide_content = []
            
            # הוספת סגנון רקע
            background_style = self._get_background_style(slide)
            
            # יצירת מיכל לתוכן השקופית
            slide_content.append(f'<div class="slide-content" style="{background_style}">')
            
            # המרת כל הצורות בשקופית
            for shape_index, shape in enumerate(slide.shapes):
                try:
                    element_id = f"slide{index}_shape{shape_index}"
                    animation_style = self._add_animation_to_element(element_id, shape_index, index)
                    
                    if shape.has_text_frame:
                        # המרת מסגרת טקסט
                        shape_html = self._convert_text_frame(shape, element_id, animation_style)
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # המרת תמונה
                        shape_html = self._convert_picture(shape, element_id, animation_style)
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # המרת טבלה
                        shape_html = self.table_handler.convert_table(shape, element_id, animation_style)
                    else:
                        # המרת צורה אחרת
                        shape_html = self._convert_shape(shape, element_id, animation_style)
                    
                    if shape_html:
                        slide_content.append(shape_html)
                except Exception as e:
                    logger.error(f"שגיאה בהמרת צורה {shape_index} בשקופית {index}: {str(e)}")
                    continue
            
            slide_content.append('</div>')
            
            return {
                'content': '\n'.join(slide_content)
            }
        except Exception as e:
            logger.error(f"שגיאה בהמרת שקופית {index}: {str(e)}")
            return {
                'content': f'<div class="error">שגיאה בהמרת שקופית {index + 1}</div>'
            }

    def _get_background_style(self, slide):
        """מחזיר את סגנון הרקע של השקופית"""
        try:
            background = slide.background
            if background.fill.type:
                if hasattr(background.fill, 'solid') and background.fill.solid:
                    color = background.fill.fore_color
                    if hasattr(color, 'rgb'):
                        rgb = color.rgb
                        return f"background-color: rgb({rgb[0]}, {rgb[1]}, {rgb[2]});"
                elif hasattr(background.fill, 'pattern') and background.fill.pattern:
                    # טיפול בדפוסי רקע
                    fore_color = background.fill.fore_color
                    back_color = background.fill.back_color
                    if hasattr(fore_color, 'rgb') and hasattr(back_color, 'rgb'):
                        fore_rgb = fore_color.rgb
                        back_rgb = back_color.rgb
                        return f"""
                            background-color: rgb({back_rgb[0]}, {back_rgb[1]}, {back_rgb[2]});
                            background-image: linear-gradient(45deg, 
                                rgb({fore_rgb[0]}, {fore_rgb[1]}, {fore_rgb[2]}) 25%, 
                                transparent 25%, transparent 75%, 
                                rgb({fore_rgb[0]}, {fore_rgb[1]}, {fore_rgb[2]}) 75%);
                            background-size: 10px 10px;
                        """
        except Exception as e:
            logger.debug(f"שגיאה בקבלת רקע השקופית: {str(e)}")
        
        # ברירת מחדל - רקע לבן
        return "background-color: white;"

    def _convert_text_frame(self, shape, element_id="", animation_style=""):
        """ממיר מסגרת טקסט ל-HTML"""
        if not shape.text.strip():
            return ""
            
        text_frame = shape.text_frame
        paragraphs_html = []
        
        for paragraph in text_frame.paragraphs:
            if not paragraph.text.strip():
                continue
                
            p_style = self._get_paragraph_style(paragraph)
            p_html = f'<p style="{p_style}">'
            
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                    
                run_style = self._get_run_style(run)
                run_html = f'<span style="{run_style}">{run.text}</span>'
                p_html += run_html
                
            p_html += '</p>'
            paragraphs_html.append(p_html)
            
        if not paragraphs_html:
            return ""
            
        return f'<div id="{element_id}" style="position: absolute; {self._get_shape_position(shape)}; {animation_style}">\n' + \
               '\n'.join(paragraphs_html) + \
               '</div>'

    def _get_paragraph_style(self, paragraph):
        """מחזיר את הסגנון של פסקה"""
        styles = []
        
        # יישור טקסט
        alignment = {
            0: 'right',  # RIGHT
            1: 'left',   # LEFT
            2: 'center', # CENTER
            3: 'justify' # JUSTIFY
        }
        styles.append(f"text-align: {alignment.get(paragraph.alignment, 'right')}")
        
        # מרווח בין שורות
        if paragraph.line_spacing:
            styles.append(f"line-height: {paragraph.line_spacing}")
            
        # מרווח לפני ואחרי
        if paragraph.space_before:
            styles.append(f"margin-top: {paragraph.space_before.pt}pt")
        if paragraph.space_after:
            styles.append(f"margin-bottom: {paragraph.space_after.pt}pt")
            
        return '; '.join(styles)

    def _get_run_style(self, run):
        """מחזיר את הסגנון של קטע טקסט"""
        styles = []
        
        # גופן
        if run.font.name:
            styles.append(f"font-family: '{run.font.name}'")
            
        # גודל
        if run.font.size:
            styles.append(f"font-size: {run.font.size.pt}pt")
            
        # צבע
        if run.font.color.rgb:
            rgb = run.font.color.rgb
            styles.append(f"color: rgb({rgb[0]}, {rgb[1]}, {rgb[2]})")
            
        # מודגש
        if run.font.bold:
            styles.append("font-weight: bold")
            
        # נטוי
        if run.font.italic:
            styles.append("font-style: italic")
            
        # קו תחתון
        if run.font.underline:
            styles.append("text-decoration: underline")
            
        return '; '.join(styles)

    def _get_shape_position(self, shape):
        """מחזיר את מיקום הצורה"""
        left = shape.left.pt
        top = shape.top.pt
        width = shape.width.pt
        height = shape.height.pt
        
        return f"left: {left}pt; top: {top}pt; width: {width}pt; height: {height}pt"

    def _convert_picture(self, shape, element_id="", animation_style=""):
        """ממיר תמונה ל-HTML"""
        try:
            # המרת מיקום וגודל מEMU ל-פיקסלים
            left = shape.left / 914400 * 96
            top = shape.top / 914400 * 96
            width = shape.width / 914400 * 96
            height = shape.height / 914400 * 96
            
            # המרת התמונה ל-base64
            image = shape.image
            image_bytes = image.blob
            image_base64 = base64.b64encode(image_bytes).decode()
            content_type = image.content_type or 'image/png'
            
            style = f"""
                position: absolute;
                left: {left:.2f}px;
                top: {top:.2f}px;
                width: {width:.2f}px;
                height: {height:.2f}px;
                object-fit: contain;
                {animation_style}
            """
            
            return f'''
                <img 
                    class="image" 
                    id="{element_id}" 
                    src="data:{content_type};base64,{image_base64}"
                    style="{style}"
                    alt="תמונה במצגת"
                />
            '''
        except Exception as e:
            print(f"שגיאה בהמרת תמונה: {str(e)}")
            return ''

    def _convert_shape(self, shape, element_id="", animation_style=""):
        """ממיר צורה ל-HTML"""
        # המרת מיקום וגודל מEMU ל-פיקסלים
        left = shape.left / 914400 * 96
        top = shape.top / 914400 * 96
        width = shape.width / 914400 * 96
        height = shape.height / 914400 * 96
        
        # חילוץ צבע הרקע
        fill_color = "transparent"
        if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
            color = shape.fill.fore_color
            if hasattr(color, 'rgb'):
                rgb = color.rgb
                fill_color = f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
        
        # חילוץ צבע המסגרת
        border_color = "none"
        border_width = "0"
        if shape.line:
            if hasattr(shape.line, 'color') and shape.line.color:
                color = shape.line.color
                if hasattr(color, 'rgb'):
                    rgb = color.rgb
                    border_color = f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
            if hasattr(shape.line, 'width'):
                border_width = f"{shape.line.width / 914400 * 96:.2f}px"
        
        style = f"""
            position: absolute;
            left: {left:.2f}px;
            top: {top:.2f}px;
            width: {width:.2f}px;
            height: {height:.2f}px;
            background-color: {fill_color};
            border: {border_width} solid {border_color};
            border-radius: {self._get_shape_radius(shape)};
            {animation_style}
        """
        
        return f'<div class="shape" id="{element_id}" style="{style}"></div>'

    def _get_shape_radius(self, shape):
        """מחזיר את רדיוס הפינות של הצורה"""
        try:
            # אם זה מלבן מעוגל
            if hasattr(shape, 'adjustment_values') and shape.adjustment_values:
                return min(shape.width, shape.height) * shape.adjustment_values[0] / 100000
        except:
            pass
        return 0

    def generate_embed_code(self, html_path, width="100%", height="600px", responsive=True):
        """יוצר קוד להטמעה עבור Wix"""
        try:
            # קריאת תוכן ה-HTML
            with open(html_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # יצירת גרסה מותאמת להטמעה
            embed_content = self._modify_for_embed(content)
            
            # שמירת הגרסה להטמעה
            embed_path = html_path.replace('.html', '_embed.html')
            with open(embed_path, 'w', encoding='utf-8') as f:
                f.write(embed_content)
            
            # יצירת קוד ההטמעה
            iframe_code = f'''
<!-- קוד להטמעת המצגת -->
<div class="presentation-container" style="position: relative; width: {width};">
    <iframe 
        src="{os.path.basename(embed_path)}"
        style="width: 100%; height: {height}; border: none; overflow: hidden;"
        allowfullscreen="true"
        loading="lazy"
    ></iframe>
</div>

<!-- סקריפט להתאמה אוטומטית של גובה -->
<script>
window.addEventListener('message', function(e) {{
    if (e.data && e.data.type === 'resize') {{
        const iframe = document.querySelector('.presentation-container iframe');
        if (iframe) {{
            iframe.style.height = e.data.height + 'px';
        }}
    }}
}});
</script>
'''
            
            # שמירת קוד ההטמעה לקובץ
            embed_code_path = html_path.replace('.html', '_embed_code.txt')
            with open(embed_code_path, 'w', encoding='utf-8') as f:
                f.write(iframe_code)
            
            return embed_code_path, embed_path
            
        except Exception as e:
            logger.error(f"שגיאה ביצירת קוד להטמעה: {str(e)}")
            raise

    def _modify_for_embed(self, content):
        """מתאים את ה-HTML להטמעה"""
        # הוספת קוד JavaScript להתאמת גובה
        resize_script = """
<script>
function updateParentHeight() {{
    const height = document.documentElement.scrollHeight;
    window.parent.postMessage({{
        type: 'resize',
        height: height
    }}, '*');
}}

// עדכון בטעינה ובשינוי גודל
window.addEventListener('load', updateParentHeight);
window.addEventListener('resize', updateParentHeight);

// עדכון בשינוי שקופית
document.addEventListener('DOMContentLoaded', () => {{
    const slides = document.querySelectorAll('.slide');
    const observer = new MutationObserver(updateParentHeight);
    
    slides.forEach(slide => {{
        observer.observe(slide, {{
            attributes: true,
            attributeFilter: ['style']
        }});
    }});
}});
</script>
"""
        # הוספת הסקריפט לפני סגירת תג body
        content = content.replace('</body>', f'{resize_script}</body>')
        
        # הוספת מטא תגים להטמעה בטוחה
        meta_tags = """
    <meta http-equiv="X-Frame-Options" content="SAMEORIGIN">
    <meta http-equiv="Content-Security-Policy" content="frame-ancestors 'self' *.wix.com">
"""
        content = content.replace('</head>', f'{meta_tags}</head>')
        
        return content
