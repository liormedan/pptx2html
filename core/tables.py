from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import _Cell

class TableHandler:
    @staticmethod
    def get_table_css():
        """מחזיר את הגדרות ה-CSS לטבלאות"""
        return """
            .pptx-table {
                border-collapse: collapse;
                width: 100%;
                height: 100%;
                position: absolute;
                background: white;
            }
            
            .pptx-table td {
                border: 1px solid #ccc;
                padding: 8px;
                position: relative;
                vertical-align: middle;
                box-sizing: border-box;
            }
            
            .pptx-table tr:nth-child(even) {
                background-color: #f9f9f9;
            }
            
            .pptx-table tr:hover {
                background-color: #f5f5f5;
            }
            
            @media (max-width: 768px) {
                .pptx-table {
                    font-size: 14px;
                }
                
                .pptx-table td {
                    padding: 4px;
                }
            }
        """

    def convert_table(self, shape, element_id="", animation_style=""):
        """ממיר טבלת PowerPoint ל-HTML"""
        if not hasattr(shape, 'table'):
            return ''
            
        table = shape.table
        
        # חישוב המיקום והגודל
        left = shape.left if hasattr(shape, 'left') else 0
        top = shape.top if hasattr(shape, 'top') else 0
        width = shape.width if hasattr(shape, 'width') else 'auto'
        height = shape.height if hasattr(shape, 'height') else 'auto'
        
        style = f"""
            left: {left}px;
            top: {top}px;
            width: {width}px;
            height: {height}px;
            {animation_style}
        """
        
        # יצירת הטבלה
        html = f'<div class="table-container" id="{element_id}" style="{style}">'
        html += '<table class="pptx-table">'
        
        # מעבר על כל השורות והתאים
        for row in table.rows:
            html += '<tr>'
            for cell in row.cells:
                # עיצוב התא
                cell_style = self._get_cell_style(cell)
                html += f'<td style="{cell_style}">'
                
                # המרת תוכן התא
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        para_style = self._get_paragraph_style(paragraph)
                        html += f'<p style="{para_style}">{paragraph.text}</p>'
                
                html += '</td>'
            html += '</tr>'
            
        html += '</table></div>'
        return html
        
    def _get_cell_style(self, cell):
        """מחזיר את הסגנון של תא בטבלה"""
        style_parts = []
        
        # צבע רקע
        if cell._tc.get_or_add_tcPr().fill:
            fill = cell._tc.get_or_add_tcPr().fill
            if hasattr(fill, 'fore_color') and fill.fore_color:
                rgb = fill.fore_color.rgb
                if rgb:
                    style_parts.append(f"background-color: rgb({rgb[0]}, {rgb[1]}, {rgb[2]})")
        
        # גבולות
        if cell._tc.get_or_add_tcPr().borders:
            borders = cell._tc.get_or_add_tcPr().borders
            for border in ['top', 'right', 'bottom', 'left']:
                border_obj = getattr(borders, border)
                if border_obj:
                    width = border_obj.width
                    if hasattr(border_obj, 'color') and border_obj.color:
                        rgb = border_obj.color.rgb
                        if rgb:
                            style_parts.append(
                                f"border-{border}: {width}pt solid rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
                            )
        
        # יישור טקסט
        if cell._tc.get_or_add_tcPr().vertical_anchor:
            valign_map = {
                'top': 'flex-start',
                'middle': 'center',
                'bottom': 'flex-end'
            }
            valign = valign_map.get(cell._tc.get_or_add_tcPr().vertical_anchor, 'center')
            style_parts.append(f"vertical-align: {valign}")
        
        return "; ".join(style_parts)
        
    def _get_paragraph_style(self, paragraph):
        """מחזיר את הסגנון של פסקה בתא"""
        style_parts = ["margin: 0"]
        
        # יישור טקסט
        if paragraph.alignment is not None:
            alignment_map = {
                0: 'right',  # PP_ALIGN.LEFT in RTL
                1: 'center',
                2: 'left',   # PP_ALIGN.RIGHT in RTL
                3: 'justify'
            }
            style_parts.append(f"text-align: {alignment_map.get(paragraph.alignment, 'right')}")
            
        return "; ".join(style_parts)
